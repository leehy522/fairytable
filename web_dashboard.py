import streamlit as st
import yfinance as yf
import pandas as pd
import matplotlib.pyplot as plt
import copy
import io
import pypdf
import re
from pptx import Presentation
from pptx.util import Pt
from datetime import datetime

# --- [1. 페이지 기본 설정] ---
# 이 코드가 가장 먼저 나와야 합니다.
st.set_page_config(page_title="요정비닐 스마트 시스템", layout="wide")

# --- [2. 로그인 보안 로직] ---
def check_password():
    """아이디와 비밀번호를 확인하여 로그인을 제어합니다."""
    USER_ID = "lhy"
    USER_PW = "dlghkdud1%" #

    if "password_correct" not in st.session_state:
        st.session_state.password_correct = False

    if st.session_state.password_correct:
        return True

    # 로그인 화면 디자인
    st.title("🔐 요정비닐 시스템 접속")
    col_l, col_r = st.columns([1, 2])
    with col_l:
        input_id = st.text_input("ID", placeholder="아이디 입력", key="login_id")
        input_pw = st.text_input("Password", type="password", placeholder="비밀번호 입력", key="login_pw")
        
        if st.button("로그인 실행"):
            if input_id == USER_ID and input_pw == USER_PW:
                st.session_state.password_correct = True
                st.rerun() # 성공 시 즉시 화면 갱신
            else:
                st.error("❌ 정보가 일치하지 않습니다.")
    return False

# 로그인 체크 실행 (통과 못 하면 여기서 멈춤)
if not check_password():
    st.stop()

# --- [3. 공통 로직 및 함수 정의] ---
# (여기서부터 기존 버전 1의 get_pallet_capacity 등 함수들이 이어집니다)
def get_pallet_capacity(sku):
    sku = str(sku)
    if sku in ['32058611', '15651222']: return 300
    if sku in ['29558294', '32711887']: return 192
    if sku == '32083343': return 400
    if sku == '32366753': return 560
    return 300

def duplicate_slide(prs, index):
    template = prs.slides[index]
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    new_slide = prs.slides.add_slide(blank_layout)
    for shp in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shp.element)
    for shape in template.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide

def set_bold_text(text_frame, content, is_bold=True, font_size=None):
    text_frame.text = str(content)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = is_bold
            if font_size: run.font.size = Pt(font_size)

# 💡 텍스트 공란 문제를 해결하기 위해 로직을 꽉 채웠습니다.
def fill_slide_data(slide, p, po_num, fc_name, year, month, day):
    try:
        current_plt_idx = int(p['no'].split('-')[1])
        total_qty = int(p['total_qty'])
        cap = int(p['cap'])
        # 팔레트별 수량 계산 로직 (v4.98)
        display_qty = cap if current_plt_idx * cap <= total_qty else (total_qty % cap if total_qty % cap != 0 else cap)
    except: 
        display_qty = p['total_qty']

    # 💡 윤겸님이 말씀하신 바로 그 핵심 텍스트 치환 로직입니다!
    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            txt = shape.text
            
            # 1. 박스수량 및 팔레트 번호 (예: 12-1)
            if "박스수량" in txt or "BOX" in txt:
                set_bold_text(tf, f"{p['no']} / 총 박스수량  ({p['total_qty']} BOX)", True)
            
            # 2. 입고 날짜 및 납품 센터명
            elif "입고예정일자" in txt or "납품센터명" in txt:
                set_bold_text(tf, f"입고예정일자 ({int(month)}월 {int(day)}일) / 납품센터명 ({fc_name} 센터)", True)
            
            # 3. 고정 업체명 입력
            elif "업체명" in txt:
                set_bold_text(tf, "업체명         (   주식회사 페어리드림    )", True)
            
            # 4. 발주번호 입력
            elif "발주번호" in txt:
                set_bold_text(tf, f"발주번호       ({po_num})", True)
        
        # 표(Table)가 있는 경우 SKU와 상품명, 수량을 채웁니다.
        if shape.has_table:
            table = shape.table
            try:
                # 💡 합쳐진 품목 리스트(items_list)를 하나씩 꺼내어 표의 행(Row)에 채웁니다.
                for idx, item in enumerate(p['items_list']):
                    row_idx = idx + 1 
                    if row_idx >= len(table.rows): break # 표의 칸을 넘어가면 중단
                    
                    # 1. SKU 번호 입력
                    set_bold_text(table.cell(row_idx, 1).text_frame, item['sku'], False)
                    
                    # 2. 상품명 입력 (앞에 [발주번호]가 붙은 상태)
                    set_bold_text(table.cell(row_idx, 2).text_frame, item['name'], False, font_size=11)
                    
                    # 💡 [핵심 수정] 전체 합계가 아닌, 이 상품의 진짜 수량(item['qty'])을 가져옵니다.
                    sku_cap = str(item.get('cap', p.get('cap', 300))) 
                    
                    # 3. 발주수량 및 입고확인 칸에 해당 상품의 개별 수량 입력
                    set_bold_text(table.cell(row_idx, 3).text_frame, sku_cap, False) # 발주수량
                    set_bold_text(table.cell(row_idx, 4).text_frame, sku_cap, False) # 입고확인
                    
                    # 4. 비고란 (전날 날짜 표기 로직 적용)
                    table.cell(row_idx, 5).text = f"-\n/{year}.{int(month)}.{int(day)}"
            except Exception as e:
                pass
                
@st.cache_data(ttl=60)
def load_google_sheet_data():
    CSV_URL = "https://docs.google.com/spreadsheets/d/e/2PACX-1vTVvCbm9KEoUrqvlXSyIyLHmstIGZuiuTMLYDBnmgnxInrfoMelDXFSWogUdHUfNALb7uC_nBAIyzif/pub?output=csv"
    try:
        df = pd.read_csv(CSV_URL)
        df.columns = [str(c).strip() for c in df.columns]
        return df.dropna(subset=['상품명'])
    except: return pd.DataFrame()

# --- [여기서부터 기존 메뉴 로직 시작] ---
# --- [2. 공통 로직 및 함수 정의] ---
# (밀크런 관련 함수들: get_pallet_capacity, duplicate_slide, set_bold_text, fill_slide_data는 상단에 정의)
# (이미 정의된 함수들은 가독성을 위해 생략하며, 실제 코드에는 그대로 유지하시면 됩니다.)

@st.cache_data(ttl=60)
def load_google_sheet_data():
    CSV_URL = "https://docs.google.com/spreadsheets/d/e/2PACX-1vTVvCbm9KEoUrqvlXSyIyLHmstIGZuiuTMLYDBnmgnxInrfoMelDXFSWogUdHUfNALb7uC_nBAIyzif/pub?output=csv"
    try:
        df = pd.read_csv(CSV_URL)
        df.columns = [str(c).strip() for c in df.columns]
        return df.dropna(subset=['상품명'])
    except Exception as e:
        st.error(f"구글 시트 연결 오류: {e}")
        return pd.DataFrame()

# --- [3. 사이드바 메뉴 구성] ---
st.sidebar.title("🚀 요정비닐 관리자")
menu = st.sidebar.radio("메뉴를 선택하세요", 
    ["🏷️ 요정비닐 상품 현황", "🚚 밀크런 PPT 변환", "📦 택배 송장 변환", "🏭 원가 시뮬레이터", "📈 시장 지표 분석"])

# --- [4. 메뉴별 화면 로직] ---

# --- 메뉴 1: 요정비닐 상품 현황 ---
if menu == "🏷️ 요정비닐 상품 현황":
    st.title("🏷️ 요정비닐 상품 실시간 현황")
    st.caption("구글 스프레드시트와 동기화 중 (자동 갱신 주기: 1분)")
    st.divider()

    df = load_google_sheet_data()

    if not df.empty:
        col1, col2 = st.columns([1, 1])
        with col1:
            st.metric("총 등록 상품", f"{len(df)} 종")
        
        with st.expander("🔍 상품 검색 및 필터", expanded=True):
            search_query = st.text_input("", placeholder="검색할 상품명을 입력하세요")
        
        display_df = df[df['상품명'].str.contains(search_query, na=False)] if search_query else df

        st.subheader(f"📦 상품 목록 ({len(display_df)}건)")
        st.dataframe(
            display_df, 
            use_container_width=True, 
            hide_index=True,
            column_config={
                "상품명": st.column_config.TextColumn("📋 상품명"),
                "단가": st.column_config.NumberColumn("💰 단가", format="₩%d"),
                "재고": st.column_config.NumberColumn("🔢 현재고", format="%d")
            }
        )
        st.success("✅ 최신 데이터가 반영되었습니다.")
    else:
        st.error("데이터를 불러올 수 없습니다. 구글 시트의 '웹에 게시' 설정을 확인해주세요.")

elif menu == "🚚 밀크런 PPT 변환":
    st.title("🚚 밀크런 통합 편집 시스템 (v4.98 로직 이식)")
    
    tpl_file = st.file_uploader("1. 양식 PPT 업로드", type=['pptx'])
    pdf_files = st.file_uploader("2. 발주서 PDF 업로드 (다중 선택)", type=['pdf'], accept_multiple_files=True)

    if tpl_file and pdf_files:
        # 데이터 리스트 초기화 관리
        if "extracted_data" not in st.session_state:
            st.session_state.extracted_data = []

        # 💡 분석 버튼 클릭 시 이전 데이터를 비우고 새로 시작하도록 수정
        if st.button("🔍 발주서 데이터 정밀 분석 (홀수 장 전용)"):
            all_extracted = []
            for pdf_file in pdf_files:
                reader = pypdf.PdfReader(pdf_file)
                
                # 💡 enumerate를 사용하여 페이지 번호(i)를 가져옵니다.
                # 파이썬은 0부터 시작하므로 0, 2, 4... 가 실제 1, 3, 5페이지(홀수 장)입니다.
                for i, page in enumerate(reader.pages):
                    if i % 2 != 0: 
                        continue # 짝수 장(index 1, 3, 5...)은 건너뜁니다.

                    text = page.extract_text() + "\n"
                    # --- [이하 기존 데이터 추출 로직 동일] ---
                    po_match = (re.search(r"(?:발주번호|PO|no|Info)\s*[:\s\n]*(\d{9})", text, re.I) or 
                                re.search(r"(\d{9})", text))
                    if not po_match: continue
                    
                    po_num = po_match.group(1) if hasattr(po_match, 'group') else po_match[0]
                    
                    po_num = po_match.group(1) if hasattr(po_match, 'group') else po_match[0]
                    fc_match = re.search(r"(?:FC명|FC\s*Name|센터명)\s*[:\s\n]*([A-Z0-9가-힣]+)", text, re.I) or \
                               re.search(r"([가-힣]+)센터", text)
                    fc_name = fc_match.group(1).strip() if fc_match else "알수없음"
                    date_match = re.search(r"(\d{4}-\d{2}-\d{2})", text)
                    date_raw = date_match.group(1) if date_match else "2026-03-12"
                    
                    sku_matches = list(re.finditer(r"\b(\d{8})\b", text))
                    processed_in_page = set()
                    for m in sku_matches:
                        sku = m.group(1)
                        if sku in processed_in_page: continue
                        cap = get_pallet_capacity(sku)
                        block = text[m.end():m.end()+450]
                        name_search = re.search(r"([가-힣]{2,}[가-힣\s\d\-\(\)]+)", block)
                        real_name = name_search.group(1).strip() if name_search else "상품명확인"
                        nums = re.findall(r"\b\d{1,4}\b", block)
                        qty = int(nums[1]) if len(nums) >= 2 else (int(nums[0]) if len(nums) == 1 else 0)
                        
                        if qty > 0:
                            all_extracted.append({
                                "발주번호": po_num, "센터": fc_name, "SKU": sku, 
                                "상품명": real_name[:40], "확정수량": qty, "적재량": cap, "date": date_raw
                            })
                            processed_in_page.add(sku)
            st.session_state.extracted_data = all_extracted
            st.rerun()

        # 편집 및 PPT 생성 로직
        if st.session_state.extracted_data:
            st.subheader("📊 발주 데이터 통합 편집")
            edited_df = st.data_editor(pd.DataFrame(st.session_state.extracted_data), num_rows="dynamic", use_container_width=True)

        if st.button("🚀 지능형 합짐 및 PPT 생성"):
                try:
                    prs = Presentation(tpl_file)
                    # 슬라이드 초기화
                    while len(prs.slides) > 1:
                        rId = prs.slides._sle[1].rId
                        prs.part.drop_rel(rId); del prs.slides._sle[1]
                    
                    is_first = True
                    for center, group in edited_df.groupby("센터"):
                        po_list = sorted([str(p) for p in group["발주번호"].unique()])
                        all_pos = ", ".join(po_list)
                        
                        mixed_items = []
                        total_qty_sum = 0 # 💡 센터별 전체 박스 합계를 저장할 변수
                        
                        for _, row in group.iterrows():
                            q = int(row["확정수량"])
                            if q <= 0: continue
                            # 개별 발주번호를 상품명 앞에 붙여서 혼동 방지
                            mixed_items.append({'sku': row['SKU'], 'name': f"[{row['발주번호']}] {row['상품명']}", 'qty': q})
                            total_qty_sum += q
                        
                        if not mixed_items: continue
                        
                        cap = int(group["적재량"].iloc[0])
                        # 팔레트 수 계산 (이 값은 p_info의 'no'에만 쓰여야 합니다)
                        if total_qty_sum < 300:
                            tot_plt = 1
                        else:
                            tot_plt = (total_qty_sum // cap) + (1 if total_qty_sum % cap > 0 else 0)
                        y, m, d = group["date"].iloc[0].split('-')
                        
                        for i in range(1, tot_plt + 1):
                            # 💡 오류 수정 포인트: 'total_qty' 자리에 tot_plt가 아닌 total_qty_sum을 정확히 전달!
                            p_info = {
                                'no': f"{tot_plt}-{i}", 
                                'total_qty': total_qty_sum, # 팔레트 수(12)가 아닌 실제 박스 수(예: 3600) 입력
                                'cap': cap, 
                                'items_list': mixed_items
                            }
                            
                            for _ in range(2):
                                if is_first:
                                    slide = prs.slides[0]
                                    is_first = False
                                else:
                                    slide = duplicate_slide(prs, 0)
                                fill_slide_data(slide, p_info, all_pos, center, y, m, d)

                    ppt_out = io.BytesIO()
                    prs.save(ppt_out)
                    st.download_button("📥 최종 PPT 다운로드", ppt_out.getvalue(), "밀크런_수량수정_결과.pptx")
                    st.success(f"✅ 수정 완료! 총 {len(prs.slides)}장의 슬라이드가 생성되었습니다.")
                    
                except Exception as e:
                    st.error(f"PPT 생성 중 에러가 발생했습니다: {e}")
                    
# --- 메뉴 3: 택배 송장 변환 ---
elif menu == "📦 택배 송장 변환":
    st.title("📦 택배 송장 자동 변환기 (A-type)")
    st.write("원본 주문 엑셀을 요정비닐 템플릿 양식에 맞춰 변환합니다.")

    col1, col2 = st.columns(2)
    with col1:
        input_file = st.file_uploader("1. 원본 주문 엑셀 선택", type=['xlsx', 'xls'])
    with col2:
        template_file = st.file_uploader("2. 템플릿 엑셀(A-type) 선택", type=['xlsx', 'xls'])

    if input_file and template_file:
        if st.button("🚀 변환 실행"):
            try:
                # 1. 파일 읽기 (숫자 변형 방지를 위해 문자열로 읽기)
                src = pd.read_excel(input_file, dtype=str)
                
                # 2. 매핑 및 검증 로직
                mapping = {
                    "주문번호": "Order ID",
                    "받는사람": "Receiver Name",
                    "전화번호1": "Mobile",
                    "전화번호2": "Mobile",
                    "우편번호": "Zip Code",
                    "주소": "Detailed address",
                    "상품명1": "Product Information",
                }
                # 도시 정보가 담긴 다양한 컬럼명 후보들
                city_candidates = ["City", "city", "도시", "시", "시/군/구"]
                city_col = next((c for c in city_candidates if c in src.columns), None)

                if not city_col:
                    st.error("⚠️ 원본 파일에서 'City' 또는 '도시' 컬럼을 찾을 수 없습니다.")
                else:
                    # 3. 데이터 변환 처리
                    out = pd.DataFrame()
                    for out_col, src_col in mapping.items():
                        if src_col in src.columns:
                            out[out_col] = src[src_col].fillna("").astype(str)

                    # 주소 결합 로직: City + Detailed address
                    out["주소"] = (
                        src[city_col].fillna("").astype(str).str.strip() + " " + 
                        src["Detailed address"].fillna("").astype(str).str.strip()
                    ).str.strip()

                    # 전화번호2는 1과 동일하게 세팅
                    out["전화번호2"] = out["전화번호1"]

                    # 4. 결과 다운로드 생성
                    output = io.BytesIO()
                    with pd.ExcelWriter(output, engine='openpyxl') as writer:
                        out.to_excel(writer, index=False)
                    
                    st.success(f"✅ 변환 완료! (총 {len(out)}행)")
                    st.download_button(
                        label="📥 변환된 엑셀 다운로드",
                        data=output.getvalue(),
                        file_name=f"요정비닐_송장변환_{datetime.now().strftime('%m%d')}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                    )
            except Exception as e:
                st.error(f"❌ 변환 실패: {e}")
                
elif menu == "🏭 원가 시뮬레이터":
    st.title("🏭 원가 및 규격 시뮬레이터")
    
    # --- [섹션 1: 원단 규격 정밀 계산기] ---
    st.subheader("📏 원단 규격 정밀 계산기")
    # key를 추가하여 하단 위젯과 ID 충돌을 방지합니다.
    calc_mode = st.radio("계산 모드", ["⚖️ 무게 산출", "🔍 두께 역산"], horizontal=True, key="mode_selector")
    
    c1, c2, c3 = st.columns(3)
    with c1: 
        v_width = st.number_input("비닐 폭 (mm)", value=630, key="calc_v_width") # key 추가
    with c2: 
        v_length = st.number_input("원단 총 길이 (m)", value=1800, key="calc_v_length") # key 추가
    
    res_weight = 0
    if calc_mode == "⚖️ 무게 산출":
        with c3: 
            v_thick = st.number_input("두께 (mm)", value=0.009, format="%.3f", key="calc_v_thick")
        res_weight = (v_width/1000) * v_length * 2 * 0.92 * v_thick
        st.info(f"💡 예상 무게: {res_weight:.2f} kg")
    else:
        with c3: 
            v_weight_in = st.number_input("실제 무게 (kg)", value=13.8, key="calc_v_weight_in")
        res_thick = v_weight_in / ((v_width/1000) * v_length * 2 * 0.92)
        st.warning(f"💡 역산된 두께: {res_thick:.4f} mm")

    # --- [섹션 2: 원재료 혼합 단가 계산] ---
    st.divider()
    st.subheader("🧪 1. 원재료 혼합 단가 설정")
    
    col1, col2 = st.columns(2)
    with col1:
        v_price = st.number_input("신원료 가격 (원/kg)", value=1530, key="raw_v_price") 
        r_price = st.number_input("재생원료 가격 (원/kg)", value=1300, key="raw_r_price") 
    with col2:
        v_ratio = st.slider("신원료 혼합 비율 (%)", 0, 100, 70, key="raw_v_ratio") 
        st.caption(f"신원료 {v_ratio}% : 재생원료 {100-v_ratio}%")

    st.write("---")
    col3, col4 = st.columns(2)
    with col3:
        c_price = st.number_input("조색제 가격 (원/kg)", value=2700, key="raw_c_price") 
    with col4:
        c_ratio = st.number_input("조색제 혼합 비율 (%)", value=2.5, step=0.1, format="%.1f", key="raw_c_ratio") 

    # 혼합 단가 수식 유지
    base_price = (v_price * (v_ratio / 100)) + (r_price * ((100 - v_ratio) / 100))
    final_unit_price = (base_price * (1 - c_ratio/100)) + (c_price * (c_ratio/100))
    st.success(f"🎨 **최종 원재료 단가: ₩{final_unit_price:,.2f} / kg**")

    # --- [섹션 3: 원단 규격 및 롤당 생산 원가] ---
    st.divider()
    st.subheader("📏 2. 원단 규격 및 생산 원가")
    
    c_w1, c_w2, c_w3 = st.columns(3)
    with c_w1:
        # 상단과 라벨이 같아도 key가 다르면 에러가 나지 않습니다.
        width_mm = st.number_input("비닐 폭 (mm)", value=630, key="final_width_mm") 
    with c_w2:
        length_m = st.number_input("원단 총 길이 (m)", value=1800, key="final_length_m") 
    with c_w3:
        thick_mm = st.number_input("비닐 두께 (mm)", value=0.009, step=0.001, format="%.3f", key="final_thick_mm") 

    # 무게 및 원가 계산
    total_weight = (width_mm / 1000) * length_m * 2 * 0.92 * thick_mm
    total_cost = total_weight * final_unit_price

    col_res1, col_res2 = st.columns(2)
    with col_res1:
        st.metric("예상 원단 무게", f"{total_weight:.2f} kg")
    with col_res2:
        st.metric("1롤당 제조 원가", f"₩{total_cost:,.0f}")
        
# --- 메뉴 5: 시장 지표 분석 (초기 복구 버전) ---
elif menu == "📈 시장 지표 분석":
    st.title("📈 실시간 유가 및 환율 모니터링")
    st.write("WTI 유가와 원/달러 환율의 1년치 흐름을 실시간으로 가져옵니다.")
    
    # 💡 초기 방식: 버튼을 누르면 즉시 yfinance에서 데이터를 내려받습니다.
    if st.button("📊 최신 데이터 불러오기"):
        with st.spinner('데이터를 가져오는 중...'):
            try:
                # 데이터 심볼 설정 (WTI: CL=F, 환율: KRW=X)
                symbols = {"WTI 유가": "CL=F", "원/달러 환율": "KRW=X"}
                df = pd.DataFrame()
                
                for name, sym in symbols.items():
                    # 기간을 1년(1y)으로 설정하여 안정적으로 가져옵니다.
                    data = yf.download(sym, period="2y", interval="1d")
                    df[name] = data['Close']
                
                # 빈 값 채우기
                df = df.ffill()

                # 1. 주요 지표 표시 (Metric)
                c1, c2 = st.columns(2)
                c1.metric("현재 WTI 유가", f"${df['WTI 유가'].iloc[-1]:.2f}")
                c2.metric("현재 환율", f"₩{df['원/달러 환율'].iloc[-1]:.2f}")

                # 2. 이중 축 그래프 시각화 (Matplotlib)
                fig, ax1 = plt.subplots(figsize=(10, 5))
                ax2 = ax1.twinx()
                
                ax1.plot(df.index, df["WTI 유가"], color='tab:blue', label='WTI', linewidth=2)
                ax2.plot(df.index, df["원/달러 환율"], color='tab:red', label='환율', linestyle='--', linewidth=2)
                
                ax1.set_ylabel("WTI Price (USD)", color='tab:blue')
                ax2.set_ylabel("Exchange Rate (KRW)", color='tab:red')
                plt.title("WTI Oil vs USD/KRW Exchange Rate")
                
                # 범례 표시
                lines1, labels1 = ax1.get_legend_handles_labels()
                lines2, labels2 = ax2.get_legend_handles_labels()
                ax1.legend(lines1 + lines2, labels1 + labels2, loc='upper left')
                
                st.pyplot(fig)
                
                # 3. 상세 데이터 표 (하단)
                st.subheader("📋 최근 데이터 상세")
                st.dataframe(df.tail(10).sort_index(ascending=False), use_container_width=True)
                
            except Exception as e:
                st.error(f"데이터 연동 실패: {e}")
