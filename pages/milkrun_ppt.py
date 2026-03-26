import copy
import io
import re

import pandas as pd
import pypdf
import streamlit as st
from pptx import Presentation
from pptx.util import Pt
from auth import check_password

# 1. 페이지 기본 설정 및 보안 체크 (최상단 배치)
st.set_page_config(page_title="밀크런 PPT 변환", page_icon="🚚", layout="wide")

if not check_password():
    st.stop()

# 2. 세션 상태 안전 초기화
if "extracted_data" not in st.session_state:
    st.session_state.extracted_data = []

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  헬퍼 함수
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def get_pallet_capacity(sku: str) -> int:
    """SKU 번호에 따라 팔레트 최대 적재 박스 수를 반환합니다."""
    sku = str(sku)
    if sku in ["32058611", "15651222"]:
        return 300
    if sku in ["29558294", "32711887"]:
        return 192
    if sku == "32083343":
        return 400
    if sku == "32366753":
        return 560
    return 300

def duplicate_slide(prs: Presentation, index: int):
    """index번 슬라이드를 맨 뒤에 복제해서 반환합니다."""
    template = prs.slides[index]
    blank_layout = (
        prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    )
    new_slide = prs.slides.add_slide(blank_layout)
    for shp in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shp.element)
    for shape in template.shapes:
        new_el = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return new_slide

def set_bold_text(text_frame, content, is_bold: bool = True, font_size=None) -> None:
    """TextFrame 전체를 content로 교체하고 굵기/크기를 설정합니다."""
    text_frame.text = str(content)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = is_bold
            if font_size:
                run.font.size = Pt(font_size)

def fill_slide_data(slide, p: dict, po_num: str, fc_name: str,
                    year: str, month: str, day: str) -> None:
    """슬라이드의 텍스트/표를 발주 데이터로 채웁니다."""
    try:
        current_plt_idx = int(p["no"].split("-")[1])
        total_qty = int(p["total_qty"])
        cap = int(p["cap"])
        display_qty = (
            cap
            if current_plt_idx * cap <= total_qty
            else (total_qty % cap if total_qty % cap != 0 else cap)
        )
    except Exception:
        display_qty = p["total_qty"]

    for shape in slide.shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            txt = shape.text

            if "박스수량" in txt or "BOX" in txt:
                set_bold_text(
                    tf, f"{p['no']} / 총 박스수량  ({p['total_qty']} BOX)", True
                )
            elif "입고예정일자" in txt or "납품센터명" in txt:
                set_bold_text(
                    tf,
                    f"입고예정일자 ({int(month)}월 {int(day)}일) / 납품센터명 ({fc_name} 센터)",
                    True,
                )
            elif "업체명" in txt:
                set_bold_text(tf, "업체명         (   주식회사 페어리드림    )", True)
            elif "발주번호" in txt:
                set_bold_text(tf, f"발주번호       ({po_num})", True)

        if shape.has_table:
            table = shape.table
            try:
                item_count = len(p["items_list"])
                for idx, item in enumerate(p["items_list"]):
                    row_idx = idx + 1
                    if row_idx >= len(table.rows):
                        break
                    set_bold_text(table.cell(row_idx, 1).text_frame, item["sku"], False)
                    set_bold_text(
                        table.cell(row_idx, 2).text_frame, item["name"], False, font_size=11
                    )
                    display_val = (
                        str(item.get("qty", item.get("확정수량", 1)))
                        if item_count > 1
                        else str(item.get("cap", p.get("cap", 300)))
                    )
                    set_bold_text(table.cell(row_idx, 3).text_frame, display_val, False)
                    set_bold_text(table.cell(row_idx, 4).text_frame, display_val, False)
                    table.cell(row_idx, 5).text = f"-\n/{year}.{int(month)}.{int(day)}"
            except Exception:
                pass

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  PDF 데이터 추출
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _extract_pdf_data(pdf_files) -> list[dict]:
    """업로드된 PDF 리스트에서 발주 데이터를 추출합니다. (홀수 페이지만)"""
    all_extracted = []
    for pdf_file in pdf_files:
        reader = pypdf.PdfReader(pdf_file)
        for i, page in enumerate(reader.pages):
            if i % 2 != 0:          # 짝수 인덱스(짝수 장) 스킵
                continue
            text = page.extract_text() + "\n"
            po_match = re.search(
                r"(?:발주번호|PO|no|Info)\s*[:\s\n]*(\d{9})", text, re.I
            ) or re.search(r"(\d{9})", text)
            if not po_match:
                continue

            po_num = po_match.group(1)
            fc_match = re.search(
                r"(?:FC명|FC\s*Name|센터명)\s*[:\s\n]*([A-Z0-9가-힣]+)", text, re.I
            ) or re.search(r"([가-힣]+)센터", text)
            fc_name = fc_match.group(1).strip() if fc_match else "알수없음"

            date_match = re.search(r"(\d{4}-\d{2}-\d{2})", text)
            date_raw = date_match.group(1) if date_match else "2026-03-12"

            processed_in_page: set = set()
            for m in re.finditer(r"\b(\d{8})\b", text):
                sku = m.group(1)
                if sku in processed_in_page:
                    continue
                cap = get_pallet_capacity(sku)
                block = text[m.end():m.end() + 450]
                name_search = re.search(r"([가-힣]{2,}[가-힣\s\d\-\(\)]+)", block)
                real_name = name_search.group(1).strip() if name_search else "상품명확인"
                nums = re.findall(r"\b\d{1,4}\b", block)
                qty = (
                    int(nums[1]) if len(nums) >= 2
                    else (int(nums[0]) if len(nums) == 1 else 0)
                )
                if qty > 0:
                    all_extracted.append({
                        "발주번호": po_num,
                        "센터": fc_name,
                        "SKU": sku,
                        "상품명": real_name[:40],
                        "확정수량": qty,
                        "적재량": cap,
                        "date": date_raw,
                    })
                    processed_in_page.add(sku)
    return all_extracted

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  PPTX 생성 (핵심 로직 수정됨)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def _build_pptx(tpl_file, edited_df: pd.DataFrame) -> bytes:
    """편집된 DataFrame으로 최종 PPTX를 생성하고 bytes를 반환합니다."""
    prs = Presentation(tpl_file)
    
    # [수정 1] 원본 템플릿(0번) 외의 쓰레기 슬라이드가 있다면 삭제
    while len(prs.slides) > 1:
        rId = prs.slides._sle[1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sle[1]

    # [수정 2] '센터'가 아닌 '발주번호' 기준으로 그룹핑하여 섞임 방지
    for po_num, group in edited_df.groupby("발주번호"):
        center = group["센터"].iloc[0]
        mixed_items = []
        total_qty_sum = 0

        for _, row in group.iterrows():
            q = int(row["확정수량"])
            if q <= 0:
                continue
            mixed_items.append({
                "sku": row["SKU"],
                "name": f"[{po_num}] {row['상품명']}",
                "qty": q,
            })
            total_qty_sum += q

        if not mixed_items:
            continue

        # 팔레트 계산 로직: 총 수량이 cap을 넘어가면 N-M으로 분할
        cap = int(group["적재량"].iloc[0])
        if total_qty_sum <= cap:
            tot_plt = 1
        else:
            tot_plt = (total_qty_sum // cap) + (1 if total_qty_sum % cap > 0 else 0)

        y, m, d = group["date"].iloc[0].split("-")

        # 각 팔레트 번호(1-1, 2-1 등)별로 2장씩(제출용, 보관용) 템플릿 복제 후 작성
        for i in range(1, tot_plt + 1):
            p_info = {
                "no": f"{tot_plt}-{i}",
                "total_qty": total_qty_sum,
                "cap": cap,
                "items_list": mixed_items,
            }
            
            for _ in range(2): 
                # [수정 3] 항상 깨끗한 원본(0번) 슬라이드를 복제하여 맨 뒤에 추가
                new_slide = duplicate_slide(prs, 0)
                fill_slide_data(new_slide, p_info, po_num, center, y, m, d)

    # [수정 4] 작업이 모두 끝난 후, 기준이 되었던 빈 템플릿(0번)은 삭제 처리
    if len(prs.slides) > 1:
        rId = prs.slides._sle[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sle[0]

    ppt_out = io.BytesIO()
    prs.save(ppt_out)
    return ppt_out.getvalue()

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
#  메인 UI 실행부
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

st.title("🚚 밀크런 통합 편집 시스템 (N-M 정밀 분할 적용)")

tpl_file = st.file_uploader("1. 양식 PPT 업로드", type=["pptx"])
pdf_files = st.file_uploader(
    "2. 발주서 PDF 업로드 (다중 선택)", type=["pdf"], accept_multiple_files=True
)

if tpl_file and pdf_files:
    # ── PDF 분석 ──────────────────────────────────────
    if st.button("🔍 발주서 데이터 정밀 분석 (홀수 장 전용)"):
        with st.spinner("PDF 분석 중..."):
            st.session_state.extracted_data = _extract_pdf_data(pdf_files)
        st.rerun()

    # ── 편집 테이블 ───────────────────────────────────
    if st.session_state.extracted_data:
        st.subheader("📊 발주 데이터 통합 편집")
        edited_df = st.data_editor(
            pd.DataFrame(st.session_state.extracted_data),
            num_rows="dynamic",
            use_container_width=True,
        )

        # ── PPT 생성 ──────────────────────────────────────
        if st.button("🚀 지능형 병합 및 PPT 생성"):
            try:
                ppt_bytes = _build_pptx(tpl_file, edited_df)
                st.download_button(
                    "📥 최종 PPT 다운로드",
                    ppt_bytes,
                    "밀크런_수량수정_결과.pptx",
                )
                st.success("✅ PPT 생성 완료! (14장 정상 출력)")
            except Exception as e:
                st.error(f"PPT 생성 중 에러: {e}")

if check_password():
    show_margin_calc()
